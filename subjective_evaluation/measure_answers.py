# -*- coding: utf-8 -*-
"""
Created on Mon Mar 06 15:23:30 2017

@author: sakurai
"""

import os

import numpy as np
import matplotlib.pyplot as plt
import pptx

from make_pptx import emu_to_px

if __name__ == '__main__':
    answers_dir = 'answers'
    labels_path = 'labels.txt'

    r_true = []  # ground truth of aspect ratios
    for image_path, _ in np.loadtxt(labels_path, str):
        image = plt.imread(image_path)
        height, width = image.shape[:2]
        aspect_ratio = float(width) / height
        r_true.append(aspect_ratio)
    l_true = np.log(r_true)  # ground truth of logarithm of aspect ratios

    # measure individual errors for each image for each subject
    l_human = []
    subject_names = []
    for pptx_name in os.listdir(answers_dir):
        pptx_path = os.path.join(answers_dir, pptx_name)
        subject_name = os.path.splitext(pptx_name)[0]
        subject_names.append(subject_name)
        presentation = pptx.Presentation(pptx_path)

        l_human_i = []
        for page, slide in enumerate(presentation.slides):
            if page == 0:
                continue

            picture = slide.shapes[1]
            height = emu_to_px(picture.height)
            width = emu_to_px(picture.width)

            r = float(width) / height  # aspect ratio
            l = np.log(r)  # logarithm of aspect ratio
            l_human_i.append(l)

        l_human.append(l_human_i)
    l_human = np.array(l_human)
    l_error = l_human - l_true

    print('Erros (l_human - l_true):')
    print(l_error)
    print()

    image_numbers = np.arange(len(l_true)) + 1
    plt.plot(image_numbers, l_error.T, '.', markersize=10)
    plt.title('Errors ')
    plt.xticks(image_numbers)
    plt.xlabel('image ID')
    plt.ylabel('Error (l_human - l_true)')
    plt.legend(subject_names, bbox_to_anchor=(1, 1))
    plt.grid()

    average_absolute_error = np.absolute(l_error).mean()
    print('Average absolute error (log scale):', average_absolute_error)
    print()

    subjectwise_aae = np.absolute(l_error).mean(axis=1)
    print('Individual scores (log scale)')
    for rank, i in enumerate(np.argsort(subjectwise_aae)):
        print('{:2} {:<11} {}'.format(
            rank + 1, subject_names[i], subjectwise_aae[i]))
    print()

    print('Average absolute error: {:.2f} %'.format(
        (np.exp(average_absolute_error) - 1) * 100))
    print('Individual scores')
    for rank, i in enumerate(np.argsort(subjectwise_aae)):
        print('{:2} {:<11} {:.2f} %'.format(
            rank + 1, subject_names[i], (np.exp(subjectwise_aae[i]) - 1) * 100)
            )
    print()

    imagewise_aae = np.absolute(l_error).mean(axis=0)
    print('Image-wise average absolute errors')
    for i in np.argsort(imagewise_aae):
        print('{:2} {}'.format(i + 1, imagewise_aae[i]))
