# -*- coding: utf-8 -*-
"""
Created on Mon Feb 27 17:10:15 2017

@author: sakurai
"""

import math
import os
import platform
import subprocess
import numpy as np
import matplotlib.pyplot as plt
import pptx


def is_buildable(data_dir_path='.'):
    return not bool(subprocess.check_output('git status -s ' + data_dir_path))


def emu_to_px(emu):
    _EMUS_PER_PX = 9525.0 if platform.system() == 'Windows' else 12700.0
    return int(round(emu / _EMUS_PER_PX))


def px_to_emu(px):
    _EMUS_PER_PX = 9525 if platform.system() == 'Windows' else 12700
    return px * _EMUS_PER_PX


def stretch(picture, ratio, in_log=False):
    '''
    Args:
        picture (pptx.shapes.picture.Picture):
            Picture object to be stretched.

        ratio (float):
            Ratio of horizontal stretch. `ratio` must be non-negative float
            number if `in_log` is False.

        in_log (bool):
            If this is set to True, then `ratio` is interpreted as in
            logarithm. For example, "ratio=1, in_log=False" is equivalent to
            "ratio=0, in_log=True".
    '''
    if in_log:
        ratio = math.exp(ratio)

    picture.width = int(picture.width * ratio)


def make_pptx(images_dir_path, labels_path, output_file_path,
              in_log=True, target_long_edge_len=500):
    prs = pptx.Presentation()
    slide_layout = prs.slide_layouts[5]  # title only layout

    for i, (image_path, ratio) in enumerate(np.loadtxt(labels_path, str)):
        ratio = float(ratio)
        slide = prs.slides.add_slide(slide_layout)
        title_box = slide.shapes[0]
        title_box.text = str(i + 1)

        picture = slide.shapes.add_picture(image_path, 0, 0)
        stretch(picture, ratio, in_log=in_log)

        actua_long_edge_len = max(picture.height, picture.width)
        scale = px_to_emu(target_long_edge_len) / float(actua_long_edge_len)
        picture.height = int(picture.height * scale)
        picture.width = int(picture.width * scale)

        picture.left = 0
        picture.top = prs.slide_height - picture.height

    prs.save(output_file_path)

if __name__ == '__main__':
    images_dir_path = 'images'
    labels_path = 'labels.txt'
    output_file_path = 'task.pptx'
    in_log = True
    target_long_edge_len = 500  # in px

    make_pptx(images_dir_path, labels_path, output_file_path, in_log,
              target_long_edge_len)
