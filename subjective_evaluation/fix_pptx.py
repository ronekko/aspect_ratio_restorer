# -*- coding: utf-8 -*-
"""
Created on Thu Mar 02 15:08:55 2017

@author: yamane
"""

import os
import numpy as np
import platform
import subprocess
import matplotlib.pyplot as plt
import cv2
import pptx

from chainer import serializers

import voc2012_regression_max_pooling
import utility


def is_buildable(data_dir_path='.'):
    return not bool(subprocess.check_output('git status -s ' + data_dir_path))


def emu_to_px(emu):
    _EMUS_PER_PX = 9525.0 if platform.system() == 'Windows' else 12700.0
    return int(round(emu / _EMUS_PER_PX))


def px_to_emu(px):
    _EMUS_PER_PX = 9525 if platform.system() == 'Windows' else 12700
    return px * _EMUS_PER_PX


def make_fix_pptx(model, labels_path, output_file_path,
                  target_long_edge_len=500):
    prs = pptx.Presentation()
    slide_layout = prs.slide_layouts[5]  # title only layout

    for i, (image_path, ratio) in enumerate(np.loadtxt(labels_path, str)):
        slide = prs.slides.add_slide(slide_layout)
        title_box = slide.shapes[0]
        title_box.text = str(i + 1)

        img = plt.imread(image_path)
        dis_img = utility.change_aspect_ratio(img, np.exp(float(ratio)), 1)
        square = utility.crop_center(dis_img)
        resize = cv2.resize(square, (256, 256))
        crop = utility.crop_224(resize)
        batch = crop[None, ...]
        batch = np.transpose(batch, (0, 3, 1, 2))
        batch = batch.astype(np.float32)
        y = model.predict(batch, True)
        r = np.exp(y)
        fix = utility.change_aspect_ratio(dis_img, 1/r, 1)
        fix_image_path = os.path.join(save_path, str(i+1)) + '.jpg'
        plt.imsave(fix_image_path, fix)

        picture = slide.shapes.add_picture(fix_image_path, 0, 0)

        actua_long_edge_len = max(picture.height, picture.width)
        scale = px_to_emu(target_long_edge_len) / float(actua_long_edge_len)
        picture.height = int(picture.height * scale)
        picture.width = int(picture.width * scale)

        picture.left = 0
        picture.top = prs.slide_height - picture.height

    prs.save(output_file_path)


if __name__ == '__main__':
    labels_path = 'labels.txt'
    output_file_path = 'pptx/fix_pptx.pptx'
    target_long_edge_len = 500  # in px
    # モデルのルートパス
    model_file = r'C:\Users\yamane\Dropbox\correct_aspect_ratio\dog_data_regression_ave_pooling\1485768519.06_asp_max_4.0\dog_data_regression_ave_pooling.npz'
    # 修正画像を保存する場所
    save_path = r'C:\Users\yamane\Dropbox\correct_aspect_ratio\subjective_evaluation\fix_images'
    # モデル読み込み
    model = voc2012_regression_max_pooling.Convnet().to_gpu()
    # Optimizerの設定
    serializers.load_npz(model_file, model)

    make_fix_pptx(model, labels_path, output_file_path, target_long_edge_len)
