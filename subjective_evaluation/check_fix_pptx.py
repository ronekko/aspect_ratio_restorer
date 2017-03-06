# -*- coding: utf-8 -*-
"""
Created on Sat Mar 04 18:59:06 2017

@author: yamane
"""

import platform
import numpy as np
import pptx


def emu_to_px(emu):
    _EMUS_PER_PX = 9525.0 if platform.system() == 'Windows' else 12700.0
    return int(round(emu / _EMUS_PER_PX))


def original_asp():
    # 正常画像のパワポファイルの場所
    ori_pptx_path = r'C:\Users\yamane\Dropbox\correct_aspect_ratio\subjective_evaluation\pptx\origin_pptx.pptx'
    r_ori = []
    presentation = pptx.Presentation(ori_pptx_path)
    for slide in presentation.slides:

        picture = slide.shapes[1]
        height = emu_to_px(picture.height)
        width = emu_to_px(picture.width)

        r = float(width) / height  # aspect ratio
        r_ori.append(r)

    return r_ori


if __name__ == '__main__':
    # 修正済みのパワポファイルの置き場所
    fix_pptx_path = r'C:\Users\yamane\Dropbox\correct_aspect_ratio\subjective_evaluation\pptx\make_fix_pptx.pptx'
    r_ori = original_asp()

    r_fix = []
    l_fix = []
    presentation = pptx.Presentation(fix_pptx_path)
    i = 0
    for slide in presentation.slides:
        title = slide.shapes[0]
        print title.text

        picture = slide.shapes[1]
        height = emu_to_px(picture.height)
        width = emu_to_px(picture.width)

        r = (float(width) / height) / r_ori[i]  # aspect ratio
        i += 1
        l = np.log(r)  # logarithm of aspect ratio
        r_fix.append(r)
        l_fix.append(l)
        print 'aspect ratio =', r
        print 'logarighm of AR =', l
        print

    print 'mean error aspect ratio =', np.exp(np.mean(np.abs(l_fix)))
    print 'error logarighm of AR =', np.mean(np.abs(l_fix))
    print
