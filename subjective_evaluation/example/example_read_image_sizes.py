# -*- coding: utf-8 -*-
"""
Created on Thu Feb 23 14:04:14 2017

@author: sakurai
"""

import platform
import numpy as np
import matplotlib.pyplot as plt
import pptx


def emu_to_px(emu):
    _EMUS_PER_PX = 9525.0 if platform.system() == 'Windows' else 12700.0
    return int(round(emu / _EMUS_PER_PX))


if __name__ == '__main__':
    presentation = pptx.Presentation('example.pptx')
    for slide in presentation.slides:
        title = slide.shapes[0]
        print title.text

        picture = slide.shapes[1]
        height = emu_to_px(picture.height)
        width = emu_to_px(picture.width)
        top = emu_to_px(picture.top)
        left = emu_to_px(picture.left)
        print '(top, left) = ({}, {})'.format(top, left)
        print '(height, width) = ({}, {})'.format(height, width)

        r = float(width) / height  # aspect ratio
        l = np.log(r)  # logarithm of aspect ratio
        print 'aspect ratio =', r
        print 'logarighm of AR =', l
        print
